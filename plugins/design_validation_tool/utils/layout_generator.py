import os
import traceback
from collections import defaultdict

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

from qgis.core import (QgsFillSymbol, QgsFeatureRequest, QgsLayout,  # type: ignore
                       QgsLayoutExporter, QgsLayoutItemMap, QgsLayoutPoint,
                       QgsLayoutSize, QgsProject, QgsRectangle,
                       QgsRuleBasedRenderer, QgsSimpleFillSymbolLayer,
                       QgsUnitTypes)
from qgis.PyQt.QtCore import QRectF  # type: ignore
from qgis.PyQt.QtCore import Qt  # type: ignore
from qgis.PyQt.QtGui import QColor  # type: ignore
from qgis.utils import iface  # type: ignore

# ---------------------------------------------------------------------------
# Slide geometry constants  (widescreen 16:9)
# ---------------------------------------------------------------------------
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)

_IMG_X = Inches(0)
_IMG_Y = Inches(0)
_IMG_W = Inches(6.5)
_IMG_H = Inches(7.5)

_DIV_X = Inches(6.55)   # thin divider line x-position

_TEXT_X = Inches(6.75)
_TEXT_Y = Inches(0.35)
_TEXT_W = Inches(6.33)
_TEXT_H = Inches(6.9)

# Colors
_CLR_BG        = RGBColor(0xF7, 0xF8, 0xFA)
_CLR_PANEL_BG  = RGBColor(0xFF, 0xFF, 0xFF)
_CLR_DIVIDER   = RGBColor(0xCC, 0xCC, 0xCC)
_CLR_RULE_BG   = RGBColor(0x1A, 0x56, 0xDB)   # blue badge background
_CLR_RULE_FG   = RGBColor(0xFF, 0xFF, 0xFF)   # white badge text
_CLR_LABEL     = RGBColor(0x6B, 0x72, 0x80)   # muted label colour
_CLR_VALUE     = RGBColor(0x11, 0x18, 0x27)   # near-black value text
_CLR_NOTE_BG   = RGBColor(0xFF, 0xF3, 0xCD)   # light amber note box
_CLR_NOTE_FG   = RGBColor(0x92, 0x40, 0x09)   # dark amber note text
_CLR_HEADING   = RGBColor(0x11, 0x18, 0x27)

_STREET_FIELD_CANDIDATES = [
    'STREETNAME', 'streetname', 'street_name', 'StreetName',
    'STREET', 'street', 'NAAM', 'naam', 'straat', 'STRAAT',
    'name', 'NAME', 'label', 'LABEL',
]


# ===========================================================================
# QGIS symbol helpers
# ===========================================================================

def uncheck_locked_symbols(project, target_layer_names=None):
    original_states = {}

    if target_layer_names:
        layers_to_process = []
        for layer_name in target_layer_names:
            layers = project.mapLayersByName(layer_name)
            if layers:
                layers_to_process.extend(layers)
    else:
        layers_to_process = list(project.mapLayers().values())

    for layer in layers_to_process:
        if not hasattr(layer, 'renderer') or not layer.renderer():
            continue
        renderer = layer.renderer()
        if isinstance(renderer, QgsRuleBasedRenderer):
            original_states[layer.id()] = {}
            root_rule = renderer.rootRule()
            modified = False

            def process_rules(rule, path=""):
                nonlocal modified
                rule_label = rule.label() if rule.label() else ""
                rule_key = rule.ruleKey()
                if "Locked" in rule_label:
                    original_states[layer.id()][rule_key] = rule.active()
                    if rule.active():
                        rule.setActive(False)
                        modified = True
                for child_rule in rule.children():
                    process_rules(child_rule, path)

            for child_rule in root_rule.children():
                process_rules(child_rule)

            if modified:
                layer.setRenderer(renderer.clone())
                layer.triggerRepaint()

    if iface:
        iface.mapCanvas().refreshAllLayers()
    return original_states


def restore_locked_symbols(project, original_states):
    if not original_states:
        return
    for layer_id, rule_states in original_states.items():
        layer = project.mapLayer(layer_id)
        if not layer or not layer.renderer():
            continue
        renderer = layer.renderer()
        if not isinstance(renderer, QgsRuleBasedRenderer):
            continue
        root_rule = renderer.rootRule()
        modified = False

        def restore_rules(rule):
            nonlocal modified
            rule_key = rule.ruleKey()
            if rule_key in rule_states:
                original_state = rule_states[rule_key]
                if rule.active() != original_state:
                    rule.setActive(original_state)
                    modified = True
            for child_rule in rule.children():
                restore_rules(child_rule)

        for child_rule in root_rule.children():
            restore_rules(child_rule)

        if modified:
            layer.setRenderer(renderer.clone())
            layer.triggerRepaint()


# ===========================================================================
# Street-name lookup
# ===========================================================================

def _find_street_layer(project):
    for layer in project.mapLayers().values():
        if not layer.isValid():
            continue
        field_names = [f.name() for f in layer.fields()]
        for candidate in _STREET_FIELD_CANDIDATES:
            if candidate in field_names:
                return layer, candidate
    return None, None


def get_street_name_for_geometry(geom, project, street_layer=None, street_field=None):
    if not geom or geom.isEmpty():
        return None
    if street_layer is None or street_field is None:
        street_layer, street_field = _find_street_layer(project)
    if street_layer is None:
        return None

    centroid = geom.centroid()
    search_rect = centroid.buffer(200, 8).boundingBox()
    request = QgsFeatureRequest().setFilterRect(search_rect)
    best_dist = float('inf')
    best_name = None

    for f in street_layer.getFeatures(request):
        fgeom = f.geometry()
        if not fgeom or fgeom.isEmpty():
            continue
        dist = fgeom.distance(centroid)
        if dist < best_dist:
            val = f[street_field]
            if val is not None and str(val).strip() not in ('', 'NULL', 'None'):
                best_dist = dist
                best_name = str(val).strip()

    return best_name


# ===========================================================================
# Main entry point
# ===========================================================================

def run_report(iface, layer_name_filter, output_dir, report_title, max_violations):
    """
    Two-phase report generation.

    Phase 1 – all QGIS/Qt operations: export map screenshots, collect metadata.
    Phase 2 – all python-pptx/lxml operations: build the PowerPoint presentation.

    The two phases never interleave so that QGIS's libxml2 usage and
    python-pptx's lxml usage cannot corrupt each other's shared state.
    """
    print("=== Starting Feature Report Generation ===")
    project = QgsProject.instance()
    locked_symbol_states = {}

    try:
        locked_symbol_states = uncheck_locked_symbols(project)
        iface.mapCanvas().refreshAllLayers()

        all_layers = list(QgsProject.instance().mapLayers().values())
        matching_layers = [
            layer for layer in all_layers
            if layer_name_filter.lower() in layer.name().lower()
        ]

        if not matching_layers:
            return False

        for layer in matching_layers:
            print(f"  - {layer.name()} ({layer.featureCount()} features)")

        try:
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
        except Exception as e:
            print(f"ERROR: Cannot create/access output directory: {e}")
            return False

        canvas = iface.mapCanvas()
        if not canvas:
            return False

        street_layer, street_field = _find_street_layer(project)

        # ------------------------------------------------------------------ #
        # PHASE 1: QGIS operations – export screenshots, collect metadata     #
        # ------------------------------------------------------------------ #
        entries = []

        for layer in matching_layers:
            if not layer.isValid():
                continue

            features = list(layer.getFeatures())
            if max_violations:
                features = features[:max_violations]

            groups = defaultdict(list)
            for feature in features:
                vtype = _safe_field(feature, 'vio_type') or 'unknown'
                groups[vtype].append(feature)

            for vtype, group_features in groups.items():
                first_feature = group_features[0]
                total_cnt = _safe_int(first_feature, 'total_cnt') or len(group_features)
                collapsed = total_cnt > 5
                features_to_render = [first_feature] if collapsed else group_features

                for feature in features_to_render:
                    try:
                        feature_id = get_feature_id(feature)
                        geom = feature.geometry()

                        if not geom or geom.isEmpty():
                            print(f"WARNING: Feature {feature_id} has no geometry, skipping")
                            continue

                        bbox = geom.boundingBox()
                        if bbox.isEmpty():
                            continue

                        buffer_x = max(bbox.width() * 0.02, 10)
                        buffer_y = max(bbox.height() * 0.02, 10)
                        buffered_bbox = QgsRectangle(
                            bbox.xMinimum() - buffer_x,
                            bbox.yMinimum() - buffer_y,
                            bbox.xMaximum() + buffer_x,
                            bbox.yMaximum() + buffer_y,
                        )

                        canvas.setExtent(buffered_bbox)
                        canvas.refresh()

                        layout = create_feature_layout(project, buffered_bbox, layer.crs())
                        image_path = os.path.join(output_dir, f"feature_{feature_id}.png")

                        if not export_layout_image(layout, image_path):
                            print(f"FAILED: Could not export {image_path}")
                            continue

                        street_name = get_street_name_for_geometry(
                            geom, project, street_layer, street_field
                        )

                        entries.append({
                            'image_path': image_path,
                            'rule_id':    _safe_field(feature, 'rule_id') or 'N/A',
                            'descr':      _safe_field(feature, 'descr') or 'N/A',
                            'details':    _safe_field(feature, 'details') or 'N/A',
                            'street':     street_name,
                            'collapsed':  collapsed,
                            'total_cnt':  total_cnt,
                        })

                    except Exception as e:
                        print(f"ERROR processing feature: {e}")
                        traceback.print_exc()
                        continue

        if not entries:
            print("ERROR: No features were successfully exported")
            return False

        # ------------------------------------------------------------------ #
        # PHASE 2: python-pptx operations – build PowerPoint presentation     #
        # No QGIS/Qt calls after this point.                                  #
        # ------------------------------------------------------------------ #
        prs = Presentation()
        prs.slide_width  = _SLIDE_W
        prs.slide_height = _SLIDE_H

        blank_layout = prs.slide_layouts[6]  # completely blank

        # Cover / title slide
        _add_title_slide(prs, blank_layout, report_title, len(entries))

        for entry in entries:
            _add_violation_slide(prs, blank_layout, entry)

        pptx_path = os.path.join(output_dir, "feature_report.pptx")
        try:
            prs.save(pptx_path)
        except Exception as e:
            print(f"ERROR: Failed to save PowerPoint: {e}")
            return False

        return True

    except Exception as e:
        traceback.print_exc()
        return False
    finally:
        restore_locked_symbols(project, locked_symbol_states)
        iface.mapCanvas().refreshAllLayers()
        print("=== Feature Report Generation Completed ===")


# ===========================================================================
# Slide builders  (called only in Phase 2 – no QGIS objects)
# ===========================================================================

def _set_slide_bg(slide, color):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_slide(prs, layout, title, violation_count):
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, RGBColor(0x1A, 0x56, 0xDB))

    # Title text
    tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(36)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.8))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = f"{violation_count} violation{'s' if violation_count != 1 else ''} found"
    run2.font.size = Pt(18)
    run2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)


def _add_violation_slide(prs, layout, entry):
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, _CLR_BG)

    # ---- Left panel: map screenshot ----
    if os.path.exists(entry['image_path']):
        slide.shapes.add_picture(
            entry['image_path'],
            _IMG_X, _IMG_Y,
            _IMG_W, _IMG_H,
        )

    # Vertical divider line
    _add_divider(slide)

    # ---- Right panel: white card ----
    _add_right_panel_bg(slide)

    # ---- Text content ----
    tb = slide.shapes.add_textbox(_TEXT_X, _TEXT_Y, _TEXT_W, _TEXT_H)
    tf = tb.text_frame
    tf.word_wrap = True

    # Rule ID badge line
    p = tf.paragraphs[0]
    p.space_after = Pt(2)
    _run(p, entry['rule_id'], bold=True, size=Pt(9), color=_CLR_RULE_FG,
         highlight_bg=_CLR_RULE_BG)

    # Description heading
    _add_para(tf, entry['descr'],
              bold=True, size=Pt(14), color=_CLR_HEADING,
              space_before=Pt(8), space_after=Pt(14))

    # Separator
    _add_para(tf, "─" * 42, size=Pt(7), color=_CLR_DIVIDER, space_after=Pt(10))

    # Street name
    _add_label_block(tf, "Street Name", entry['street'] or "—")

    # Issue description
    _add_label_block(tf, "Issue", entry['details'], space_before=Pt(12))

    # Collapsed / pop-zone warning
    if entry['collapsed']:
        _add_para(tf, "", size=Pt(6), space_before=Pt(0), space_after=Pt(0))
        _add_para(tf, "─" * 42, size=Pt(7), color=_CLR_DIVIDER,
                  space_before=Pt(10), space_after=Pt(8))
        p_note = _add_para(tf, "Note", bold=True, size=Pt(10),
                           color=_CLR_NOTE_FG, space_before=Pt(0), space_after=Pt(3))
        _add_para(
            tf,
            f"This issue occurred {entry['total_cnt']} times in this pop zone. "
            f"Please check the entire pop zone for this violation.",
            size=Pt(10), color=_CLR_NOTE_FG,
            space_before=Pt(0), space_after=Pt(0),
        )


def _add_label_block(tf, label, value, space_before=Pt(0), space_after=Pt(4)):
    """Two-line block: bold grey label + value text."""
    # Label
    p_lbl = _add_para(tf, label.upper(), bold=True, size=Pt(8),
                      color=_CLR_LABEL, space_before=space_before, space_after=Pt(2))
    # Value
    _add_para(tf, value or "—", size=Pt(11), color=_CLR_VALUE,
              space_before=Pt(0), space_after=space_after)


def _add_para(tf, text, bold=False, size=None, color=None,
              space_before=None, space_after=None):
    """Append a paragraph with a single uniformly-styled run."""
    p = tf.add_paragraph()
    if space_before is not None:
        p.space_before = space_before
    if space_after is not None:
        p.space_after = space_after
    run = p.add_run()
    run.text = text
    if bold:
        run.font.bold = True
    if size:
        run.font.size = size
    if color:
        run.font.color.rgb = color
    return p


def _run(p, text, bold=False, size=None, color=None, highlight_bg=None):
    """Add a run to an existing paragraph *p*."""
    run = p.add_run()
    run.text = text
    if bold:
        run.font.bold = True
    if size:
        run.font.size = size
    if color:
        run.font.color.rgb = color
    return run


def _add_divider(slide):
    """Add a thin vertical line between the two columns."""
    from pptx.util import Emu as _Emu
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT = 1
        _DIV_X, Inches(0),
        _DIV_X, _SLIDE_H,
    )
    line.line.color.rgb = _CLR_DIVIDER
    line.line.width = Pt(0.75)


def _add_right_panel_bg(slide):
    """White filled rectangle behind the right text panel."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1
        Inches(6.55), Inches(0),
        Inches(6.78), _SLIDE_H,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _CLR_PANEL_BG
    shape.line.color.rgb = _CLR_PANEL_BG


# ===========================================================================
# QGIS layout helpers
# ===========================================================================

def create_feature_layout(project, extent, crs):
    layout = QgsLayout(project)
    layout.initializeDefaults()

    try:
        if layout.pageCollection().pageCount() > 0:
            page = layout.pageCollection().page(0)
            page.setPageSize(QgsLayoutSize(210, 297, QgsUnitTypes.LayoutMillimeters))
            symbol = QgsFillSymbol()
            symbol.deleteSymbolLayer(0)
            simple_fill = QgsSimpleFillSymbolLayer()
            simple_fill.setColor(QColor(255, 255, 255))
            simple_fill.setStrokeStyle(Qt.NoPen)
            symbol.appendSymbolLayer(simple_fill)
            page.setPageStyleSymbol(symbol)
    except Exception as e:
        print(f"Warning: Could not configure page style: {e}")

    map_item = QgsLayoutItemMap(layout)
    map_item.setFrameEnabled(False)
    map_item.setBackgroundEnabled(True)
    map_item.setBackgroundColor(QColor(255, 255, 255))

    try:
        map_item.attemptSetSceneRect(QRectF(10, 10, 190, 260))
    except Exception:
        try:
            map_item.attemptMove(QgsLayoutPoint(10, 10, QgsUnitTypes.LayoutMillimeters))
            map_item.attemptResize(QgsLayoutSize(190, 260, QgsUnitTypes.LayoutMillimeters))
        except Exception as e:
            print(f"Error setting map item position/size: {e}")

    map_item.setCrs(crs)
    map_item.setExtent(extent)
    layout.addLayoutItem(map_item)
    return layout


def export_layout_image(layout, output_path):
    try:
        layout.refresh()
        exporter = QgsLayoutExporter(layout)
        settings = QgsLayoutExporter.ImageExportSettings()
        settings.dpi = 300
        settings.cropToContents = True
        settings.imageFormat = "PNG"
        result = exporter.exportToImage(output_path, settings)
        return result == QgsLayoutExporter.Success and os.path.exists(output_path)
    except Exception:
        traceback.print_exc()
        return False


# ===========================================================================
# Misc helpers
# ===========================================================================

def _safe_field(feature, field_name):
    try:
        val = feature[field_name]
        if val is None:
            return None
        s = str(val).strip()
        return s if s not in ('', 'NULL', 'None') else None
    except Exception:
        return None


def _safe_int(feature, field_name):
    try:
        val = feature[field_name]
        return int(val) if val is not None else None
    except Exception:
        return None


def get_feature_id(feature):
    id_fields = ['id', 'ID', 'fid', 'FID', 'objectid', 'OBJECTID', 'gid', 'GID']
    for field_name in id_fields:
        try:
            if field_name in [f.name() for f in feature.fields()]:
                value = feature[field_name]
                if value is not None:
                    return str(value)
        except Exception:
            continue
    return str(feature.id())
